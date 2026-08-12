import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_three_slide_with_storage():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: White Theme matching the user's slide template
    COLOR_BG = RGBColor(255, 255, 255)         # Pure White
    COLOR_CARD = RGBColor(248, 250, 252)       # Slate 50 (Very light gray)
    COLOR_BORDER_LIGHT = RGBColor(226, 232, 240) # Slate 200 (Light gray border)
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)     # Slate 900 (Dark Slate text)
    COLOR_TEXT_MUTED = RGBColor(71, 85, 105)   # Slate 600 (Muted text)
    COLOR_BLUE = RGBColor(29, 78, 216)         # Blue 700 (Title and accents)
    COLOR_AMBER = RGBColor(217, 119, 6)        # Amber 600 (Warnings/Before)
    COLOR_GREEN = RGBColor(22, 163, 74)        # Green 600 (Clean status/After)

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_slide_header(slide, title_text, slide_number_text=""):
        # Header container (simulating the user's template container style)
        header_container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(0.4), Inches(7.333), Inches(0.8)
        )
        header_container.fill.solid()
        header_container.fill.fore_color.rgb = COLOR_BLUE
        header_container.line.color.rgb = COLOR_BLUE
        
        tf = header_container.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        if slide_number_text:
            p.text = slide_number_text + "  |  "
            p.font.name = 'Calibri'
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_AMBER
            
            run = p.add_run()
            run.text = title_text
            run.font.name = 'Calibri'
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        else:
            p.text = title_text
            p.font.name = 'Calibri'
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

    slide_layout = prs.slide_layouts[6] # blank layout

    # ==========================================
    # SLIDE 1: Hiện trạng dữ liệu thô & Quy trình xử lý
    # ==========================================
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)
    add_slide_header(slide1, "HIỆN TRẠNG DỮ LIỆU & QUY TRÌNH XỬ LÝ", "SLIDE 1")

    # Left Side: Insert Missing Data Rate Chart Image
    chart_path = r"d:\Profile\Visual Code File\DATN_DP02_NYC\data\data clean\missing_data_chart.png"
    if os.path.exists(chart_path):
        slide1.shapes.add_picture(chart_path, Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    else:
        # Fallback card
        card_l = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
        card_l.fill.solid()
        card_l.fill.fore_color.rgb = COLOR_CARD
        card_l.line.color.rgb = COLOR_BORDER_LIGHT
        tf_l = card_l.text_frame
        p = tf_l.paragraphs[0]
        p.text = "[CHART IMAGE]"

    # Right Column: Raw data state & Issues Card
    card_r = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_BORDER_LIGHT
    card_r.line.width = Pt(1)

    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.3)

    p = tf_r.paragraphs[0]
    p.text = "HIỆN TRẠNG DỮ LIỆU THÔ"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    issues = [
        ("Quy mô gốc", "49,041 dòng | 36 cột. Nguồn: NYC DOF & U.S. Census."),
        ("Khuyết thiếu", "Xảy ra diện rộng với 68,093 ô trống (6 cột khuyết >50%)."),
        ("Sai kiểu số", "Giá bán dạng int64, ngày bán dạng object, zip dạng float."),
        ("Lệch định dạng", "Địa chỉ và Quận huyện viết hoa/thường không đồng nhất."),
        ("Ngoại lệ", "Chứa giao dịch dưới $1,000 và trên $100M gây lệch phân phối.")
    ]

    for title, desc in issues:
        p_is = tf_r.add_paragraph()
        p_is.text = f"• {title}: "
        p_is.font.name = 'Calibri'
        p_is.font.size = Pt(12.5)
        p_is.font.bold = True
        p_is.font.color.rgb = COLOR_TEXT_MAIN
        p_is.space_before = Pt(12)
        
        run = p_is.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 2: Kết quả chuẩn hóa & Đối sánh Trước/Sau
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_slide_header(slide2, "DỮ LIỆU SAU KHI ĐƯỢC CHUẨN HÓA", "SLIDE 2")

    # Split screen: Left (Before) and Right (After)
    card_before = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.4))
    card_before.fill.solid()
    card_before.fill.fore_color.rgb = COLOR_CARD
    card_before.line.color.rgb = COLOR_AMBER
    card_before.line.width = Pt(1.5)

    tf_bef = card_before.text_frame
    tf_bef.word_wrap = True
    tf_bef.margin_left = tf_bef.margin_right = Inches(0.4)
    tf_bef.margin_top = Inches(0.3)

    p = tf_bef.paragraphs[0]
    p.text = "TRẠNG THÁI CŨ (BEFORE)"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    bef_points = [
        "Hơn 68,000 ô dữ liệu bị bỏ trống",
        "Kiểu dữ liệu và định dạng không đồng nhất",
        "Chữ thường/hoa viết lẫn lộn ở Borough và Address",
        "Giao dịch ngoại lệ gây sai lệch phân phối"
    ]
    for pt in bef_points:
        p_pt = tf_bef.add_paragraph()
        p_pt.text = f"✗ {pt}"
        p_pt.font.name = 'Calibri'
        p_pt.font.size = Pt(12.5)
        p_pt.font.color.rgb = COLOR_TEXT_MUTED
        p_pt.space_before = Pt(10)

    card_after = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(3.4))
    card_after.fill.solid()
    card_after.fill.fore_color.rgb = COLOR_CARD
    card_after.line.color.rgb = COLOR_GREEN
    card_after.line.width = Pt(1.5)

    tf_aft = card_after.text_frame
    tf_aft.word_wrap = True
    tf_aft.margin_left = tf_aft.margin_right = Inches(0.4)
    tf_aft.margin_top = Inches(0.3)

    p = tf_aft.paragraphs[0]
    p.text = "TRẠNG THÁI MỚI (AFTER)"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    aft_points = [
        "Điền khuyết 100% bằng Trung vị (Median) & mặc định",
        "Đồng bộ kiểu float64, datetime64[ns], string, Int64",
        "Borough viết hoa đồng bộ; Address chuyển Title Case",
        "Phân loại bằng cờ: 1,299 giao dịch nội bộ và 351 siêu sang"
    ]
    for pt in aft_points:
        p_pt = tf_aft.add_paragraph()
        p_pt.text = f"✔️ {pt}"
        p_pt.font.name = 'Calibri'
        p_pt.font.size = Pt(12.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(10)

    # 3 KPI Cards at the bottom
    kpi_details = [
        ("49,041", "DÒNG BẢO TOÀN", Inches(0.8)),
        ("6", "CỘT NHIỄU BỊ LOẠI BỎ", Inches(4.8)),
        ("68,093", "GIÁ TRỊ KHUYẾT ĐÃ ĐIỀN", Inches(8.8))
    ]

    for val, label, left_pos in kpi_details:
        kpi_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(5.4), Inches(3.7 if "GIÁ TRỊ" in label else 3.6), Inches(1.1))
        kpi_card.fill.solid()
        kpi_card.fill.fore_color.rgb = COLOR_CARD
        kpi_card.line.color.rgb = COLOR_BLUE
        kpi_card.line.width = Pt(1)

        tf_k = kpi_card.text_frame
        tf_k.word_wrap = True
        tf_k.margin_left = tf_k.margin_right = Inches(0.2)
        tf_k.margin_top = Inches(0.15)

        p_val = tf_k.paragraphs[0]
        p_val.text = val
        p_val.font.name = 'Calibri'
        p_val.font.size = Pt(22)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_GREEN
        p_val.alignment = PP_ALIGN.CENTER

        p_lbl = tf_k.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = 'Calibri'
        p_lbl.font.size = Pt(10.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_before = Pt(2)

    # Slide Footer message
    tx_footer = slide2.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.4))
    tf_foot = tx_footer.text_frame
    tf_foot.word_wrap = True
    
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "“Dữ liệu sạch và đồng nhất giúp kết quả phân tích chính xác, đáng tin cậy hơn.”"
    p_foot.font.name = 'Calibri'
    p_foot.font.size = Pt(13)
    p_foot.font.bold = True
    p_foot.font.italic = True
    p_foot.font.color.rgb = COLOR_AMBER
    p_foot.alignment = PP_ALIGN.CENTER


    # ==========================================
    # SLIDE 3: Tổ chức & Lưu trữ dữ liệu đầu ra
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_slide_header(slide3, "TỔ CHỨC & LƯU TRỮ DỮ LIỆU ĐẦU RA", "SLIDE 3")

    # Left Column: Storage Environments
    card_l3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.4))
    card_l3.fill.solid()
    card_l3.fill.fore_color.rgb = COLOR_CARD
    card_l3.line.color.rgb = COLOR_BLUE
    card_l3.line.width = Pt(1.5)

    tf_l3 = card_l3.text_frame
    tf_l3.word_wrap = True
    tf_l3.margin_left = tf_l3.margin_right = Inches(0.4)
    tf_l3.margin_top = Inches(0.3)

    p = tf_l3.paragraphs[0]
    p.text = "MÔI TRƯỜNG LƯU TRỮ CHÍNH"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    storage_envs = [
        ("Cục bộ (Local Disk)", "Lưu trữ tại D:/.../DATN_DP02_NYC"),
        ("Đám mây Git (GitHub)", "Version control tại repo DATN_N2 (nhánh NghiTran)"),
        ("Đám mây Drive (Google Drive)", "Sao lưu tệp báo cáo, slide & tập dữ liệu đầu ra")
    ]

    for env_title, env_desc in storage_envs:
        p_env = tf_l3.add_paragraph()
        p_env.text = f"• {env_title}:\n"
        p_env.font.name = 'Calibri'
        p_env.font.size = Pt(12.5)
        p_env.font.bold = True
        p_env.font.color.rgb = COLOR_TEXT_MAIN
        p_env.space_before = Pt(8)

        run = p_env.add_run()
        run.text = f"  {env_desc}"
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # Right Column: Directory tree structure
    card_r3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(3.4))
    card_r3.fill.solid()
    card_r3.fill.fore_color.rgb = COLOR_CARD
    card_r3.line.color.rgb = COLOR_BLUE
    card_r3.line.width = Pt(1.5)

    tf_r3 = card_r3.text_frame
    tf_r3.word_wrap = True
    tf_r3.margin_left = tf_r3.margin_right = Inches(0.4)
    tf_r3.margin_top = Inches(0.3)

    p = tf_r3.paragraphs[0]
    p.text = "CẤU TRÚC THƯ MỤC DỰ ÁN"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    dir_tree = [
        "DATN_DP02_NYC/ (Thư mục gốc)",
        "  ├── data/ (Thư mục dữ liệu)",
        "  │    ├── Data crawl/ (Thô: Crawl_data_NYC.csv)",
        "  │    └── data clean/ (Sạch: Dulieu_Cleaned.csv)",
        "  └── src/ (Mã nguồn làm sạch: main.py)"
    ]

    for line in dir_tree:
        p_tree = tf_r3.add_paragraph()
        p_tree.text = line
        p_tree.font.name = 'Consolas'
        p_tree.font.size = Pt(11)
        p_tree.font.color.rgb = COLOR_TEXT_MAIN if "/" in line else COLOR_TEXT_MUTED
        p_tree.space_before = Pt(4)

    # 3 Storage KPI cards
    kpi_details_s3 = [
        ("17.5 MB", "DUNG LƯỢNG FILE THÔ", Inches(0.8)),
        ("11.6 MB", "DUNG LƯỢNG FILE SẠCH", Inches(4.8)),
        ("Git & Drive", "NỀN TẢNG CLOUD BACKUP", Inches(8.8))
    ]

    for val, label, left_pos in kpi_details_s3:
        kpi_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(5.4), Inches(3.7 if "CLOUD" in label else 3.6), Inches(1.1))
        kpi_card.fill.solid()
        kpi_card.fill.fore_color.rgb = COLOR_CARD
        kpi_card.line.color.rgb = COLOR_BLUE
        kpi_card.line.width = Pt(1)

        tf_k = kpi_card.text_frame
        tf_k.word_wrap = True
        tf_k.margin_left = tf_k.margin_right = Inches(0.2)
        tf_k.margin_top = Inches(0.15)

        p_val = tf_k.paragraphs[0]
        p_val.text = val
        p_val.font.name = 'Calibri'
        p_val.font.size = Pt(22)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_GREEN
        p_val.alignment = PP_ALIGN.CENTER

        p_lbl = tf_k.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = 'Calibri'
        p_lbl.font.size = Pt(10.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_before = Pt(2)

    # Slide Footer message
    tx_footer_s3 = slide3.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.4))
    tf_foot_s3 = tx_footer_s3.text_frame
    tf_foot_s3.word_wrap = True
    
    p_foot_s3 = tf_foot_s3.paragraphs[0]
    p_foot_s3.text = "“Dữ liệu được tổ chức khoa học giúp tối ưu hóa không gian lưu trữ và dễ dàng chia sẻ.”"
    p_foot_s3.font.name = 'Calibri'
    p_foot_s3.font.size = Pt(13)
    p_foot_s3.font.bold = True
    p_foot_s3.font.italic = True
    p_foot_s3.font.color.rgb = COLOR_AMBER
    p_foot_s3.alignment = PP_ALIGN.CENTER

    # Save presentation
    os.makedirs(r"d:\Profile\Visual Code File\DATN_DP02_NYC\data\data clean", exist_ok=True)
    out_path = r"d:\Profile\Visual Code File\DATN_DP02_NYC\data\data clean\Bao_cao_Lam_sach_Du_lieu_v2.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    create_three_slide_with_storage()
