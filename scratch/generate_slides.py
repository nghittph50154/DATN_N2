import os
import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Slate Dark Theme
    COLOR_BG = RGBColor(15, 23, 42)          # Slate 900 (Dark blue/gray)
    COLOR_CARD = RGBColor(30, 41, 59)        # Slate 800 (Card background)
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # Slate 50 (White/off-white)
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400 (Gray)
    COLOR_CYAN = RGBColor(56, 189, 248)       # Cyan 400 (Highlight, accents)
    COLOR_AMBER = RGBColor(245, 158, 11)      # Amber 500 (Warnings/Actions)
    COLOR_GREEN = RGBColor(34, 197, 94)       # Green 500 (Success/Results)

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_slide_header(slide, title_text, step_text=""):
        # Header text box
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        if step_text:
            p.text = step_text.upper() + "  |  "
            p.font.name = 'Calibri'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_AMBER
            
            run = p.add_run()
            run.text = title_text
            run.font.name = 'Calibri'
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = COLOR_TEXT_MAIN
        else:
            p.text = title_text
            p.font.name = 'Calibri'
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN

        # Add horizontal line below header
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CYAN
        shape.line.color.rgb = COLOR_CYAN

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide_layout = prs.slide_layouts[6] # blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)

    # Large Main Title
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BÁO CÁO LÀM SẠCH & CHUẨN HÓA DỮ LIỆU"
    p.font.name = 'Calibri'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "Quy trình tiền xử lý dữ liệu cho dự án Phân tích Bất động sản New York (NYC)"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_before = Pt(15)

    # Stats block in the title slide
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.3), Inches(11.333), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_CYAN
    card.line.width = Pt(1.5)

    tf_card = card.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = Inches(0.4)
    tf_card.margin_top = Inches(0.3)
    
    p_card = tf_card.paragraphs[0]
    p_card.text = "QUY MÔ DỮ LIỆU BAN ĐẦU:"
    p_card.font.name = 'Calibri'
    p_card.font.size = Pt(12)
    p_card.font.bold = True
    p_card.font.color.rgb = COLOR_TEXT_MUTED
    
    p_card_stats = tf_card.add_paragraph()
    p_card_stats.text = "49,041 dòng | 36 cột   ➔   49,041 dòng | 32 cột dữ liệu chuẩn hóa"
    p_card_stats.font.name = 'Calibri'
    p_card_stats.font.size = Pt(22)
    p_card_stats.font.bold = True
    p_card_stats.font.color.rgb = COLOR_GREEN
    p_card_stats.space_before = Pt(10)

    p_card_sub = tf_card.add_paragraph()
    p_card_sub.text = "Dựa trên cleaning_log.txt  •  Không làm mất mát bất kỳ mẫu giao dịch nào (0 dòng bị xóa)"
    p_card_sub.font.name = 'Calibri'
    p_card_sub.font.size = Pt(12)
    p_card_sub.font.italic = True
    p_card_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_card_sub.space_before = Pt(10)

    # ==========================================
    # SLIDE 2: Workflow
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_slide_header(slide2, "Tổng Quan Quy Trình 6 Bước")

    # Draw 6 blocks for the 6 steps
    step_width = Inches(3.6)
    step_height = Inches(2.2)
    gaps_x = Inches(0.4)
    gaps_y = Inches(0.3)
    start_x = Inches(0.8)
    start_y = Inches(1.8)

    steps_info = [
        ("BƯỚC 1", "Loại Bỏ Trùng Lặp", "Kiểm tra trùng lặp bản ghi hoàn toàn & trùng Business Key (Address, Sale Date, Sale Price)."),
        ("BƯỚC 2", "Xử Lý Dữ Liệu Khuyết", "Loại bỏ 6 cột khuyết >50%. Điền cột số bị lệch bằng Median. Điền text bằng 'UNKNOWN'."),
        ("BƯỚC 3", "Chuẩn Hóa Kiểu Dữ Liệu", "Ép kiểu SALE_PRICE, SALE_DATE, ZIP_CODE, và YEAR_BUILT về định dạng chuẩn để phân tích."),
        ("BƯỚC 4", "ĐONG NHẤT VĂN BẢN", "Strip khoảng trắng thừa. Chuyển chữ hoa Borough/Neighborhood. Chuyển Title Case cho Address."),
        ("BƯỚC 5", "Xử Lý Ngoại Lệ (Outlier)", "Gắn cờ (flag) giao dịch giá cực thấp (<$1,000) và giao dịch siêu sang (>$100M). Kiểm tra năm xây dựng."),
        ("BƯỚC 6", "Loại Bỏ Cột Không Đóng Góp", "Kiểm tra và xác nhận không có cột nào toàn bộ là null hoặc có độ biến động bằng không (zero variance).")
    ]

    for idx, (step_num, step_name, step_desc) in enumerate(steps_info):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (step_width + gaps_x)
        y = start_y + row * (step_height + gaps_y)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_width, step_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_CYAN if idx % 2 == 0 else COLOR_AMBER
        card.line.width = Pt(1)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.15)
        
        p = tf_card.paragraphs[0]
        p.text = step_num
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_AMBER if idx % 2 == 0 else COLOR_CYAN
        
        p2 = tf_card.add_paragraph()
        p2.text = step_name
        p2.font.name = 'Calibri'
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.space_before = Pt(4)
        
        p3 = tf_card.add_paragraph()
        p3.text = step_desc
        p3.font.name = 'Calibri'
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(6)

    # ==========================================
    # SLIDE 3: Step 1 - Duplicate Removal
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_slide_header(slide3, "Loại Bỏ Trùng Lặp (Duplicate Removal)", "Step 1")

    # Left Column: Actions and Why
    tx_left = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "HÀNH ĐỘNG XỬ LÝ"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_l.add_paragraph()
    p.text = "• Kiểm tra trùng lặp hoàn toàn:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(12)
    
    p = tf_l.add_paragraph()
    p.text = "   Đảm bảo không có bản ghi (dòng) nào bị lặp lại 100% trong cơ sở dữ liệu."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Kiểm tra trùng lặp theo Business Key:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(12)
    
    p = tf_l.add_paragraph()
    p.text = "   Tổ hợp khóa định danh giao dịch bao gồm: Địa chỉ (address), Ngày bán (sale_date), Giá bán (sale_price).\n   Mục tiêu: Đảm bảo một giao dịch thực tế không bị ghi nhận nhiều lần."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Right Column: Impact (Bình luận kết quả)
    card_r = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_GREEN
    card_r.line.width = Pt(1.5)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.4)
    
    p = tf_r.paragraphs[0]
    p.text = "KẾT QUẢ & ĐÁNH GIÁ"
    p.font.name = 'Calibri'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    
    p = tf_r.add_paragraph()
    p.text = "Không phát hiện trùng lặp"
    p.font.name = 'Calibri'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(20)
    
    p = tf_r.add_paragraph()
    p.text = "Số dòng dữ liệu được bảo toàn:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(15)
    
    p = tf_r.add_paragraph()
    p.text = "49,041 dòng   ➔   49,041 dòng"
    p.font.name = 'Calibri'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_before = Pt(5)
    
    p = tf_r.add_paragraph()
    p.text = "Nhận xét: Dữ liệu đầu vào thu thập không có hiện tượng dư thừa bản ghi hoặc lỗi ghi đè giao dịch."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(30)

    # ==========================================
    # SLIDE 4: Step 2 - Missing Values (Drop columns)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_slide_header(slide4, "Xử Lý Dữ Liệu Khuyết - Loại Bỏ Cột", "Step 2.1")

    # Left: Rule explanation
    tx_left = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "TIÊU CHÍ LOẠI BỎ CỘT"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    
    p = tf_l.add_paragraph()
    p.text = "• Nguyên tắc >50% khuyết thiếu:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Các cột có tỉ lệ trống lớn hơn 50% được coi là không đủ thông tin chất lượng để tiến hành phân tích hoặc suy diễn (imputation). Việc giữ lại có thể gây nhiễu cho mô hình học máy."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(5)

    p = tf_l.add_paragraph()
    p.text = "• Tác động đến số lượng cột:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(25)
    
    p = tf_l.add_paragraph()
    p.text = "   36 cột ban đầu   ➔   30 cột\n   (Loại bỏ hoàn toàn 6 cột không khả thi)"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    p.space_before = Pt(5)

    # Right: List of dropped columns with bars/details
    card_r = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(1.8), Inches(6.8), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_AMBER
    card_r.line.width = Pt(1)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "DANH SÁCH 6 CỘT BỊ LOẠI BỎ & TỈ LỆ KHUYẾT THIẾU"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    cols_dropped = [
        ("easement", "100.00% missing", "Cột trống hoàn toàn, không có bất kỳ dữ liệu nào."),
        ("apartment_number", "75.32% missing", "Số căn hộ khuyết phần lớn, khó định danh."),
        ("SALE PRICE PER SQFT", "52.48% missing", "Đơn giá diện tích khuyết hơn một nửa."),
        ("SQFT_PER_UNIT", "52.43% missing", "Diện tích trung bình mỗi căn hộ khuyết trên 52%."),
        ("land_sqft", "52.26% missing", "Diện tích mặt đất khuyết thiếu diện rộng."),
        ("gross_sqft", "52.26% missing", "Tổng diện tích xây dựng khuyết thiếu diện rộng.")
    ]
    
    for c_name, c_rate, c_reason in cols_dropped:
        p_col = tf_r.add_paragraph()
        p_col.text = f"• {c_name} — "
        p_col.font.name = 'Calibri'
        p_col.font.size = Pt(12)
        p_col.font.bold = True
        p_col.font.color.rgb = COLOR_CYAN
        p_col.space_before = Pt(10)
        
        run = p_col.add_run()
        run.text = c_rate
        run.font.bold = True
        run.font.color.rgb = COLOR_AMBER
        
        run_desc = p_col.add_run()
        run_desc.text = f" ({c_reason})"
        run_desc.font.bold = False
        run_desc.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 5: Step 2.2 - Missing Values Imputation
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_slide_header(slide5, "Xử Lý Dữ Liệu Khuyết - Điền Dữ Liệu", "Step 2.2")

    # Left Column: Method explanation
    tx_left = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "PHƯƠNG PHÁP SUY DIỄN (IMPUTATION)"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_l.add_paragraph()
    p.text = "• Sử dụng Trung vị (Median) cho biến số:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Dữ liệu phân phối lệch mạnh (skewed). Trung vị (Median) là chỉ số vững (robust) hơn nhiều so với Trung bình (Mean), giúp tránh bị ảnh hưởng bởi các giá trị ngoại lệ cực đoan."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Điền giá trị mặc định cho biến phân loại:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(20)
    
    p = tf_l.add_paragraph()
    p.text = "   Cột địa chỉ (address) khuyết 1 dòng được điền chữ 'UNKNOWN' để duy trì tính toàn vẹn của bản ghi thay vì loại bỏ."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Kiểm tra biến mục tiêu (SALE_PRICE):"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(20)
    
    p = tf_l.add_paragraph()
    p.text = "   Đầy đủ 100% (0 nulls), đảm bảo mục tiêu học máy sạch hoàn toàn."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GREEN

    # Right Column: Table/List of Imputations
    card_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_CYAN
    card_r.line.width = Pt(1)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "CHI TIẾT ĐIỀN GIÁ TRỊ KHUYẾT THIẾU"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    imputations = [
        ("commercial_units", "23,699 nulls", "Median = 0.0"),
        ("RESIDENTIAL_RATIO", "15,492 nulls", "Median = 1.0"),
        ("residential_units", "15,408 nulls", "Median = 1.0"),
        ("total_units", "13,478 nulls", "Median = 1.0"),
        ("zip_code", "16 nulls", "Median = 11,205.0"),
        ("address (Text)", "1 null", "Filled with 'UNKNOWN'")
    ]
    
    for col, nulls, action in imputations:
        p_imp = tf_r.add_paragraph()
        p_imp.text = f"• {col}:\n"
        p_imp.font.name = 'Calibri'
        p_imp.font.size = Pt(12)
        p_imp.font.bold = True
        p_imp.font.color.rgb = COLOR_CYAN
        p_imp.space_before = Pt(8)
        
        run = p_imp.add_run()
        run.text = f"  Số lượng khuyết: {nulls}  ➔  Xử lý: "
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED
        
        run_act = p_imp.add_run()
        run_act.text = action
        run_act.font.bold = True
        run_act.font.color.rgb = COLOR_GREEN

    # ==========================================
    # SLIDE 6: Step 3 - Data Types Fixing
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_slide_header(slide6, "Chuẩn Hóa Kiểu Dữ Liệu", "Step 3")

    # Left: Explanation
    tx_left = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Ý NGHĨA CỦA ÉP KIỂU DỮ LIỆU"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_l.add_paragraph()
    p.text = "• Đảm bảo tính toán số học chính xác:\n  Tránh lỗi kiểu dữ liệu khi đưa vào thuật toán hồi quy hoặc tính các chỉ số thống kê trung bình."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(12)
    
    p = tf_l.add_paragraph()
    p.text = "• Chuẩn hóa thời gian (Time-series):\n  Giúp trích xuất các đặc trưng thời gian (Tháng, Quý, Năm, Ngày trong tuần) dễ dàng hơn."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(15)

    p = tf_l.add_paragraph()
    p.text = "• Tránh mất mát thông tin phân loại:\n  ZIP Code bản chất là mã phân vùng địa lý, không phải là số tính toán. Nếu lưu dạng số có thể mất số 0 ở đầu (ví dụ: 07001 ➔ 7001)."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(15)

    # Right: The Changes Card
    card_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_CYAN
    card_r.line.width = Pt(1)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.4)
    
    p = tf_r.paragraphs[0]
    p.text = "BẢNG CHUYỂN ĐỔI KIỂU DỮ LIỆU"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    types_fixing = [
        ("SALE_PRICE", "int64  ➔  float64", "Đảm bảo dạng số thực cho giá trị giao dịch."),
        ("SALE_DATE", "object  ➔  datetime64[ns]", "Định dạng YYYY-MM-DD (0 dòng lỗi)."),
        ("ZIP_CODE", "float64  ➔  string", "Giữ nguyên leading zeros của mã bưu chính."),
        ("YEAR_BUILT", "float64  ➔  Int64 (Nullable)", "Chuyển sang số nguyên, xử lý năm = 0 thành NaN.")
    ]
    
    for var_name, change, comment in types_fixing:
        p_t = tf_r.add_paragraph()
        p_t.text = f"• {var_name}\n"
        p_t.font.name = 'Calibri'
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(12)
        
        run = p_t.add_run()
        run.text = f"  {change}\n"
        run.font.bold = True
        run.font.color.rgb = COLOR_GREEN
        
        run_com = p_t.add_run()
        run_com.text = f"  * {comment}"
        run_com.font.bold = False
        run_com.font.italic = True
        run_com.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 7: Step 4 - Text Consistency
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_slide_header(slide7, "Đồng Nhất Văn Bản (Text Consistency)", "Step 4")

    # Left: Explanation of why text consistency matters
    tx_left = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "TẠI SAO PHẢI CHUẨN HÓA VĂN BẢN?"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_l.add_paragraph()
    p.text = "• Loại bỏ khoảng trắng thừa (Strip):"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(12)
    
    p = tf_l.add_paragraph()
    p.text = "   Khoảng trắng ở đầu/cuối chuỗi (ví dụ: ' Manhattan ' và 'Manhattan') sẽ bị hệ thống hiểu là 2 giá trị khác nhau. Cần strip sạch sẽ."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Đồng nhất kiểu chữ (Uppercase / Title case):"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Ngăn chặn sự trùng lặp phân nhóm do nhạy chữ hoa/thường (Case Sensitivity). Chuyển toàn bộ về chữ hoa (UPPERCASE) giúp dễ dàng thống kê và lọc dữ liệu."
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Right: Results
    card_r = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_CYAN
    card_r.line.width = Pt(1)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.4)
    
    p = tf_r.paragraphs[0]
    p.text = "KẾT QUẢ ÁP DỤNG THỰC TẾ"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    text_actions = [
        ("Loại bỏ khoảng trắng thừa", "15 cột text", "0 giá trị thay đổi (đã sạch)."),
        ("Chuyển viết hoa cột 'borough'", "49,041 dòng", "Thay đổi 100% dòng dữ liệu."),
        ("Chuyển viết hoa cột 'neighborhood'", "0 thay đổi", "Neighborhood gốc đã viết hoa."),
        ("Title Case cột 'address'", "49,034 dòng", "Chuẩn hóa định dạng địa chỉ dễ đọc (Ví dụ: '123 Main St').")
    ]
    
    for title, target, impact in text_actions:
        p_act = tf_r.add_paragraph()
        p_act.text = f"• {title}:\n"
        p_act.font.name = 'Calibri'
        p_act.font.size = Pt(13)
        p_act.font.bold = True
        p_act.font.color.rgb = COLOR_CYAN
        p_act.space_before = Pt(10)
        
        run = p_act.add_run()
        run.text = f"  Phạm vi: {target}  ➔  Tác động: "
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED
        
        run_imp = p_act.add_run()
        run_imp.text = impact
        run_imp.font.bold = True
        run_imp.font.color.rgb = COLOR_GREEN

    # ==========================================
    # SLIDE 8: Step 5 & 6 - Outliers & Irrelevant Columns
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide8)
    add_slide_header(slide8, "Xử Lý Ngoại Lệ & Loại Bỏ Cột Thừa", "Steps 5 & 6")

    # Left: Outliers (Step 5)
    card_l = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_CARD
    card_l.line.color.rgb = COLOR_AMBER
    card_l.line.width = Pt(1)
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = Inches(0.4)
    tf_l.margin_top = Inches(0.3)
    
    p = tf_l.paragraphs[0]
    p.text = "BƯỚC 5: GẮN CỜ NGOẠI LỆ (OUTLIERS)"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    
    p = tf_l.add_paragraph()
    p.text = "• Nguyên tắc: Không xóa bỏ tùy tiện để bảo toàn dữ liệu thực tế, thay vào đó sử dụng kỹ thuật gắn cờ (Flagging)."
    p.font.name = 'Calibri'
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.space_before = Pt(8)
    
    p = tf_l.add_paragraph()
    p.text = "• Chuyển nhượng nội bộ (Internal Transfer):"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Giá bán < $1,000  ➔  Gắn cờ is_internal_transfer = True\n   Tác động: 1,299 dòng được gắn cờ (Thường là cho tặng gia đình)."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Giao dịch siêu sang (Luxury Deals):"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Giá bán > $100,000,000  ➔  Gắn cờ is_luxury = True\n   Tác động: 351 dòng được gắn cờ."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_l.add_paragraph()
    p.text = "• Kiểm tra năm xây dựng (YEAR_BUILT):"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_l.add_paragraph()
    p.text = "   Giới hạn 1800 - 2026. Kết quả: Tất cả các dòng hợp lệ."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_GREEN

    # Right: Irrelevant Columns (Step 6)
    card_r = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD
    card_r.line.color.rgb = COLOR_CYAN
    card_r.line.width = Pt(1)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.4)
    tf_r.margin_top = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "BƯỚC 6: LOẠI BỎ CỘT KHÔNG ĐÓNG GÓP"
    p.font.name = 'Calibri'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_r.add_paragraph()
    p.text = "• Kiểm tra cột toàn NULL (All-Null Columns):"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(20)
    
    p = tf_r.add_paragraph()
    p.text = "   Tìm các cột không chứa bất kỳ giá trị nào.\n   Kết quả: Không có cột nào toàn null sau Bước 2."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_r.add_paragraph()
    p.text = "• Kiểm tra cột không biến động (Zero Variance):"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(25)
    
    p = tf_r.add_paragraph()
    p.text = "   Tìm các cột chỉ chứa duy nhất một giá trị cho tất cả các dòng (không cung cấp thông tin phân biệt).\n   Kết quả: Không có cột nào zero variance."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 9: Summary & Conclusion
    # ==========================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide9)
    add_slide_header(slide9, "Tổng Kết Trạng Thái Dữ Liệu")

    # Big summary stats grid
    card_l = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_CARD
    card_l.line.color.rgb = COLOR_CYAN
    card_l.line.width = Pt(1.5)
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = Inches(0.4)
    tf_l.margin_top = Inches(0.4)
    
    p = tf_l.paragraphs[0]
    p.text = "CHỈ SỐ TIỀN XỬ LÝ DỮ LIỆU"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    stats_summary = [
        ("Số dòng dữ liệu ban đầu & hiện tại", "49,041", "(Bảo toàn 100% dòng)"),
        ("Số cột ban đầu", "36", ""),
        ("Số cột bị loại bỏ (>50% null)", "6", "(easement, apartment_number, land/gross sqft, etc.)"),
        ("Số cột được thêm (Gắn cờ ngoại lệ)", "2", "(is_internal_transfer, is_luxury)"),
        ("Số cột dữ liệu cuối cùng", "32", "(Tối ưu và sẵn sàng phân tích)")
    ]
    
    for label, val, sub in stats_summary:
        p_st = tf_l.add_paragraph()
        p_st.text = f"• {label}: "
        p_st.font.name = 'Calibri'
        p_st.font.size = Pt(13)
        p_st.font.color.rgb = COLOR_TEXT_MAIN
        p_st.space_before = Pt(12)
        
        run = p_st.add_run()
        run.text = val
        run.font.bold = True
        run.font.color.rgb = COLOR_GREEN
        
        if sub:
            run_sub = p_st.add_run()
            run_sub.text = f" {sub}"
            run_sub.font.italic = True
            run_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Right column: Conclusion text
    tx_right = slide9.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_r = tx_right.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "KẾT LUẬN & ĐỊNH HƯỚNG TIẾP THEO"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p = tf_r.add_paragraph()
    p.text = "1. Chất lượng dữ liệu được nâng cao:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_r.add_paragraph()
    p.text = "   Loại bỏ các cột thiếu nhiều dữ liệu và điền giá trị khuyết một cách khoa học giúp giảm sai số cho các phân tích thống kê tiếp theo."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_r.add_paragraph()
    p.text = "2. Sẵn sàng cho việc phân tích và mô hình hóa (EDA & Modeling):"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_r.add_paragraph()
    p.text = "   Dữ liệu sau khi làm sạch được lưu trữ tại tập tin Dulieu_Cleaned.csv. Các đặc trưng đã được định dạng đúng kiểu dữ liệu và gắn cờ ngoại lệ thích hợp."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p = tf_r.add_paragraph()
    p.text = "3. Triển khai dashboard:"
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(15)
    
    p = tf_r.add_paragraph()
    p.text = "   Dữ liệu sạch sẽ được kết nối trực tiếp với dashboard Streamlit để phân tích trực quan."
    p.font.name = 'Calibri'
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Ensure output directory exists
    os.makedirs(r"d:\Profile\Visual Code File\DATN_DP02_NYC\data\data clean", exist_ok=True)
    out_path = r"d:\Profile\Visual Code File\DATN_DP02_NYC\data\data clean\Bao_cao_Lam_sach_Du_lieu.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    create_presentation()
